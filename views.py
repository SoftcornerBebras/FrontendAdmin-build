from django.shortcuts import render
from ques.models import *
from ques.serializers import *
from com.models import *
from com.serializers import *
from usr.models import *
from rest_framework.views import APIView
from .serializers import *
from rest_framework.parsers import FileUploadParser,MultiPartParser
from rest_framework.decorators import parser_classes
import json
from BebrasBackend.pagination import *
from rest_framework.response import Response
from knox.auth import TokenAuthentication
from rest_framework import permissions
from BebrasBackend.constants import *
import io
import csv
import collections
import copy
import six
import os
from os import walk
import glob
#import win32com.client
#import pythoncom
from django.conf import settings
from pptx import Presentation
from datetime import datetime
import statistics
import zipfile
from operator import itemgetter


class InsertTranslation(APIView):                     #Insert translation API
    authentication_classes = (TokenAuthentication, )
    permission_classes = (permissions.IsAuthenticated,)

    def post(self,request):
        serializer = InsertTranslationSerializer(data=request.data)
        serializer.is_valid(raise_exception=True)
        serializer.save()
        return Response(status=200)

def handle_uploaded_file(f):
    with open('media/images/' + f.name, 'wb+') as destination:
        for chunk in f.chunks():
            destination.write(chunk)

@parser_classes((MultiPartParser,))
class InsertMcqQuestion(APIView):                    #Insert mcq Ques API
    authentication_classes = (TokenAuthentication, )
    permission_classes = (permissions.IsAuthenticated,)

    parser_class = (FileUploadParser,)
    def post(self,request):
        data = json.loads(request.data['data'])
        serializer = InsertMcqQuesSerializer(data=data)
        if serializer!= "":
            serializer.is_valid(raise_exception=True)
            serializer.save()
        files = request.FILES.getlist('image')
        for f in files:
            handle_uploaded_file(f)
        return Response(status=200)

@parser_classes((MultiPartParser,))
class InsertMcqWithImagesQuestion(APIView):           #Insert mcq with images Ques API
    authentication_classes = (TokenAuthentication, )
    permission_classes = (permissions.IsAuthenticated,)

    parser_class = (FileUploadParser,)
    def post(self,request):
        data = json.loads(request.data['data'])
        serializer = InsertMcqWithImagesQuesSerializer(data=data)
        if serializer!= "":
            serializer.is_valid(raise_exception=True)
            serializer.save()
        files = request.FILES.getlist('image')
        for f in files:
            handle_uploaded_file(f)
        return Response(status=200)


class ViewAgeGroupsPerQues(APIView):                   #Get Age Groups Per Ques
      authentication_classes = (TokenAuthentication, )
      permission_classes = (permissions.IsAuthenticated,)

      def get(self,request,**kwargs):
         cmpqid = QuestionAge.objects.filter(questionID=kwargs['questionID'])
         serializer = GetAgeQuestion(cmpqid, many=True)
         return Response(serializer.data)


class GetAgeGroups(APIView):          #Get Age Groups Valid for Current Year
      authentication_classes = (TokenAuthentication, )
      permission_classes = (permissions.IsAuthenticated,)

      def get(self,request):
          enddates = str(datetime.now().year)+"-12-01"
          if datetime.now().date() < datetime.strptime(enddates,'%Y-%m-%d').date():
              startrange = str(datetime.now().year-1)+"-12-01"
              lists = AgeGroup.objects.filter(created_on__range=[datetime.strptime(startrange,'%Y-%m-%d').date(),datetime.strptime(enddates,'%Y-%m-%d').date()]).order_by('AgeGroupName')
              serializer = GetAgeGroupsid(lists, many=True)
              return Response(serializer.data)
          else:
              if datetime.now().date() >= datetime.strptime(enddates,'%Y-%m-%d').date():
                   startrange = str(datetime.now().year+1)+"-12-01"
                   result = AgeGroup.objects.filter(created_on__range=[datetime.strptime(enddates,'%Y-%m-%d').date(),datetime.strptime(startrange,'%Y-%m-%d').date()]).exists()
                   if result==True:
                      lists =  AgeGroup.objects.filter(created_on__range=[datetime.strptime(enddates,'%Y-%m-%d').date(),datetime.strptime(startrange,'%Y-%m-%d').date()]).order_by('AgeGroupName')
                      serializer = GetAgeGroupsid(lists, many=True)
                      return Response(serializer.data)
                   else:
                     return Response("Redirect")

def CheckAgeGrpsForBulkUpload():
    AgeGroupIDList = []
    enddates = str(datetime.now().year) + "-12-01"
    if datetime.now().date() < datetime.strptime(enddates, '%Y-%m-%d').date():
        startrange = str(datetime.now().year - 1) + "-12-01"
        lists = AgeGroup.objects.filter(created_on__range=[datetime.strptime(startrange, '%Y-%m-%d').date(),
                                                           datetime.strptime(enddates, '%Y-%m-%d').date()])
        serializer = GetAgeGroupsid(lists, many=True)
        for i in range(0, len(serializer.data)):
            AgeGroupIDList.append(serializer.data[i]['AgeGroupName'])
            AgeGroupIDList.append(serializer.data[i]['AgeGroupID'])
        return AgeGroupIDList
    else:
        if datetime.now().date() >= datetime.strptime(enddates, '%Y-%m-%d').date():
            startrange = str(datetime.now().year + 1) + "-12-01"
            result = AgeGroup.objects.filter(created_on__range=[datetime.strptime(enddates, '%Y-%m-%d').date(),
                                                                datetime.strptime(startrange,
                                                                                  '%Y-%m-%d').date()]).exists()
            if result == True:
                lists = AgeGroup.objects.filter(created_on__range=[datetime.strptime(enddates, '%Y-%m-%d').date(),
                                                                   datetime.strptime(startrange, '%Y-%m-%d').date()])
                serializer = GetAgeGroupsid(lists, many=True)
                for i in range(0, len(serializer.data)):
                    AgeGroupIDList.append(serializer.data[i]['AgeGroupName'])
                    AgeGroupIDList.append(serializer.data[i]['AgeGroupID'])
                return AgeGroupIDList
            else:
                return "Redirect"

def InsertBulkData(reader,ser,modified_by):
    for user in reader:
        if questionTranslation.objects.filter(Identifier=user['Identifier']).exists()==False:
            countryRef = Countries.objects.get(name=user['Country'])
            domainCodeRef = code.objects.get(codeName=user['Domain'])
            quesTypeRef = code.objects.get(codeName=user['TypeOfQuestion'])
            skills = user['ComputationalSkills'].split(',')
            cs_skills = ""
            cs_skills = str(code.objects.get(codeName=skills[0]).codeID)
            for i in range(1,len(skills)):
             cs_skills = cs_skills +","+str(code.objects.get(codeName=skills[i]).codeID)
            ques = question.objects.create(countryID=countryRef, domainCodeID=domainCodeRef, questionTypeCodeID=quesTypeRef,cs_skills=cs_skills)
            langCodeRef = code.objects.get(codeName=user['Language'])
            Identifier = user['Identifier']
            modified_on = datetime.now().date()
            task = '<div style=" font-family: Arial;">'+user['Task']+'\n\n'
            Ques = '<span style="fontWeight:bold">Question: \n\n </span>'+user['Question']+'</div>'
            background = task+Ques
            explanation = '<div style=" font-family: Arial;">'+ user['Explanation'] + '</div>'
            trans={
                "caption": user['Caption'],
                "background": background,
                "explanation": explanation,
            }
            quesTrans = questionTranslation.objects.create(questionID=ques, languageCodeID=langCodeRef,
                                                           translation=trans, Identifier=Identifier, modified_by=modified_by,
                                                           modified_on=modified_on)
            AgeGroupNameList = user['AgeGroups'].split(',')
            questLevelList = user['QuestionLevels'].split(',')
            for i in range(0, len(AgeGroupNameList)):
                index = ser.index(AgeGroupNameList[i])
                AgeGroupNameRef = AgeGroup.objects.get(AgeGroupName=AgeGroupNameList[i],AgeGroupID=ser[index+1])
                quesLevelRef = code.objects.get(codeName=questLevelList[i])
                cmpQues = QuestionAge.objects.create(AgeGroupID=AgeGroupNameRef, questionID=ques,
                                                             questionLevelCodeID=quesLevelRef)
            if user['AnswerText'] == "":
                opt1 = option.objects.create(questionID=ques)
                opt2 = option.objects.create(questionID=ques)
                opt3 = option.objects.create(questionID=ques)
                opt4 = option.objects.create(questionID=ques)
                t = {
                    "option":user['Option1']
                }
                optTrans1 = optionTranslation.objects.create(optionID=opt1, languageCodeID=langCodeRef,translationO=t)
                t = {
                    "option":user['Option2']
                }
                optTrans2 = optionTranslation.objects.create(optionID=opt2, languageCodeID=langCodeRef,translationO=t)
                t = {
                    "option":user['Option3']
                }
                optTrans3 = optionTranslation.objects.create(optionID=opt3, languageCodeID=langCodeRef,translationO=t)
                t = {
                    "option":user['Option4']
                }
                optTrans4 = optionTranslation.objects.create(optionID=opt4, languageCodeID=langCodeRef,translationO=t)
                if user['CorrectOption'] == user['Option1']:
                    corrOpt = correctOption.objects.create(questionTranslationID=quesTrans, optionTranslationID=optTrans1)
                elif user['CorrectOption'] == user['Option2']:
                    corrOpt = correctOption.objects.create(questionTranslationID=quesTrans, optionTranslationID=optTrans2)
                elif user['CorrectOption'] == user['Option3']:
                    corrOpt = correctOption.objects.create(questionTranslationID=quesTrans, optionTranslationID=optTrans3)
                elif user['CorrectOption'] == user['Option4']:
                    corrOpt = correctOption.objects.create(questionTranslationID=quesTrans, optionTranslationID=optTrans4)
            else:
                  corrOpt = correctOption.objects.create(questionTranslationID=quesTrans, ansText=user['AnswerText'])
        else:
            continue


def ProcessCSVdata(reader,modified_by):
    ser = []
    lists=[]
    for user in reader:
        lists.append(user)
        if CheckAgeGrpsForBulkUpload() == "Redirect":
            return "Redirect"
        else:
            ser = CheckAgeGrpsForBulkUpload()
            AgeGroupNameList = user['AgeGroups'].split(',')
            for i in range(0, len(AgeGroupNameList)):
                if AgeGroupNameList[i] not in ser:
                    return "Redirect"
    InsertBulkData(lists,ser,modified_by)

@parser_classes((MultiPartParser,))
class BulkUploadQuestion(APIView):                                   #Bulk Upload Ques API
    authentication_classes = (TokenAuthentication, )
    permission_classes = (permissions.IsAuthenticated,)

    parser_class = (FileUploadParser,)
    def post(self, request):
        t = request.FILES['text']
        modified_by = json.loads(request.data['data'])
        f = io.TextIOWrapper(t.file)
        reader = csv.DictReader(f)
        ser = ProcessCSVdata(reader,modified_by['modified_by'])
        if ser == "Redirect":
            return Response("Redirect")
        return Response(status=200)

class InsertMarkingSchemeView(APIView):                              #Insert Marking Scheme API
     authentication_classes = (TokenAuthentication, )
     permission_classes = (permissions.IsAuthenticated,)

     def post(self, request):
        cmpAgeID = request.data.pop('competitionAgeID',None)
        cmpAgeGrpD = cmpAgeID['created_on']
        cmpAgeGrpName = cmpAgeID['AgeGroupName']
        ageID = AgeGroup.objects.get(AgeGroupName = cmpAgeGrpName, created_on=cmpAgeGrpD)
        cmpclass = AgeGroupClass.objects.filter(AgeGroupID=ageID.AgeGroupID)
        classes = []
        for i in range(0,len(cmpclass)):
            classes.append(cmpclass[i].AgeGroupClassID)
        cmpinfo = request.data.pop('CmpData',None)
        cmptype = code.objects.get(codeName=cmpinfo['cmptype'])
        cmpRef = competition.objects.get(competitionName = cmpinfo['competitionName'],startDate = cmpinfo['startDate'],competitionType = cmptype.codeID)
        cmpage = competitionAge.objects.filter(AgeGroupClassID__in = classes,competitionID=cmpRef.competitionID)
        questionLevelCodeID = request.data.pop('questionLevelCodeID', None)
        corrMarks = request.data.pop('correctMarks',None)
        incorrMarks = request.data.pop('incorrectMarks',None)
        quesLevelRef = code.objects.get(codeName= questionLevelCodeID['codeName']['name'])
        cmp_marks = ""
        for i in range(0,len(cmpage)):
            cmp_marks=competition_MarkScheme.objects.create(competitionAgeID=cmpage[i],questionLevelCodeID=quesLevelRef,correctMarks=corrMarks,incorrectMarks=incorrMarks)
        return Response(status=200)

class GetCompetition(APIView):                                        #Get All main Challenges API
     authentication_classes = (TokenAuthentication, )
     permission_classes = (permissions.IsAuthenticated,)

     def get(self,request,**kwargs):
        codeRef = code.objects.get(codeName='Main Challenge')
        lists=competition.objects.filter(competitionType = codeRef.codeID)
        serializer = CompetitionSerializer(lists, many=True)
        return Response(serializer.data)

class GetCompetitionSchoolWise(APIView):                          #Get All Cmp School Wise API
     authentication_classes = (TokenAuthentication, )
     permission_classes = (permissions.IsAuthenticated,)

     def get(self,request,**kwargs):
        schools = school.objects.get(schoolID=kwargs['schoolID'])
        codeRef = code.objects.get(codeName='Main Challenge')
        datefield = competition.objects.filter(startDate__year__gte=str(schools.registered_On)[0:4],competitionType=codeRef.codeID)
        serializer = CompetitionSerializer(datefield, many=True)
        return Response(serializer.data)

class GetNotStartedCompetitionSchoolWise(APIView):                #Get All Cmp Not Started
    authentication_classes = (TokenAuthentication, )
    permission_classes = (permissions.IsAuthenticated,)

    def get(self,request,**kwargs):
        schools = school.objects.get(schoolID=kwargs['schoolID'])
        codeRef = code.objects.get(codeName='Main Challenge')
        currentDate = datetime.now()
        datefield = competition.objects.filter(startDate__year__gte=str(schools.registered_On)[0:4],competitionType=codeRef.codeID,startDate__gte=currentDate)
        serializer = CompetitionSerializer(datefield, many=True)
        return Response(serializer.data)

class GetDistinctYears(APIView):                          #Get Distinct Years API
      authentication_classes = (TokenAuthentication, )
      permission_classes = (permissions.IsAuthenticated,)

      def get(self,request):
        lists=competition.objects.dates('startDate','year')
        return Response({"data":lists})

class GetCmpYearWise(APIView):                            #Get Cmp Year Wise API
     authentication_classes = (TokenAuthentication, )
     permission_classes = (permissions.IsAuthenticated,)

     def get(self,request,**kwargs):
        datefield = competition.objects.filter(startDate__year=kwargs['year'])
        serializer = CompetitionSerializer(datefield, many=True)
        return Response(serializer.data)

class GetSchoolClassStudents(APIView):                  #Get Students School Class Wise API
    def get(self,request,**kwargs):
        cmpid=kwargs['cmpID']
        comp= competition.objects.filter(competitionID=cmpid)
        schclass=schoolClass.objects.filter(schoolClassID=kwargs['schoolClassID'])
        agegrp=AgeGroupClass.objects.filter(ClassID=schclass[0].classNumber)
        agegroupclass=[]
        for i in range(0,len(agegrp)):
            agegroupclass.append(agegrp[i].AgeGroupClassID)
        cmpAge= competitionAge.objects.filter(AgeGroupClassID__in=agegroupclass,competitionID=comp[0].competitionID )
        lists = studentEnrollment.objects.filter(schoolClassID=kwargs['schoolClassID'],competitionAgeID=cmpAge[0].competitionAgeID).exclude(score=999)
        serializer= studentEnrollmentSerializer(lists,many=True)
        return Response(serializer.data)

class GetSchoolStudentsCmpWise(APIView):                  #Get SchoolStudent Cmp wise API
     authentication_classes = (TokenAuthentication, )
     permission_classes = (permissions.IsAuthenticated,)

     def get(self,request,**kwargs):
        cmpid = kwargs['cmpID']
        schoolclasses=schoolClass.objects.filter(schoolID=kwargs['schoolID'])
        schclasses=[]
        for i in range(0, len(schoolclasses)):
            schclasses.append(schoolclasses[i].schoolClassID)
        comp = competition.objects.filter(competitionID=cmpid)
        compAge = competitionAge.objects.filter(competitionID=comp[0].competitionID)
        cmpAgeID=[]
        for i in range(0,len(compAge)):
             cmpAgeID.append(compAge[i].competitionAgeID)
        lists = studentEnrollment.objects.filter(competitionAgeID__in=cmpAgeID,schoolClassID__in=schclasses)
        paginator = CustomPagination()
        response = paginator.generate_response(lists,studentEnrollmentSerializer,request)
        return Response(response.data)

class GetSchoolStudentsDetailCmpWise(APIView):                  #Get SchoolStudent Detail Cmp wise API
     authentication_classes = (TokenAuthentication, )
     permission_classes = (permissions.IsAuthenticated,)

     def get(self,request,**kwargs):
        cmpid = kwargs['cmpID']
        schoolclasses=schoolClass.objects.filter(schoolID=kwargs['schoolID'])
        schclasses=[]
        for i in range(0, len(schoolclasses)):
            schclasses.append(schoolclasses[i].schoolClassID)
        comp = competition.objects.filter(competitionID=cmpid)
        compAge = competitionAge.objects.filter(competitionID=comp[0].competitionID)
        cmpAgeID=[]
        for i in range(0,len(compAge)):
             cmpAgeID.append(compAge[i].competitionAgeID)
        lists = studentEnrollment.objects.filter(competitionAgeID__in=cmpAgeID,schoolClassID__in=schclasses)
        paginator = CustomPagination()
        response = paginator.generate_response(lists,studentEnrollmentSerializer,request)
        userRoleId = []
        for i in range(0,len(response.data['results'])):
            userRoleId.append(response.data['results'][i]['userID']['userID'])
            lists = UserRole.objects.filter(userID__in=userRoleId)
            serializer = UserRoleSerializer(lists,many=True)
        return Response({
            "StudData":response.data,
            "RoleData":serializer.data
        })


class DownloadSchoolStudentsCmpWise(APIView):                          #Download SchoolStudents Cmp wise
    authentication_classes = (TokenAuthentication, )
    permission_classes = (permissions.IsAuthenticated,)

    def get(self,request,**kwargs):
        cmpid = kwargs['cmpID']
        schoolclasses=schoolClass.objects.filter(schoolID=kwargs['schoolID'])
        schclasses=[]
        for i in range(0, len(schoolclasses)):
            schclasses.append(schoolclasses[i].schoolClassID)
        comp = competition.objects.filter(competitionID=cmpid)
        compAge = competitionAge.objects.filter(competitionID=comp[0].competitionID)
        cmpAgeID=[]
        for i in range(0,len(compAge)):
             cmpAgeID.append(compAge[i].competitionAgeID)
        lists = studentEnrollment.objects.filter(competitionAgeID__in=cmpAgeID,schoolClassID__in=schclasses)
        serializer= studentEnrollmentSerializer(lists,many=True)
        return Response(serializer.data)


class GetStudentsAgeGroupWise(APIView):            #Get Students AgeGrpWise API
     authentication_classes = (TokenAuthentication, )
     permission_classes = (permissions.IsAuthenticated,)

     def get(self,request,**kwargs):
        agegrpclasses=AgeGroupClass.objects.filter(AgeGroupID=kwargs['AgeID'])
        agegroupclasses=[]
        for i in range(0,len(agegrpclasses)):
            agegroupclasses.append(agegrpclasses[i].AgeGroupClassID)
        cmpage=competitionAge.objects.filter(AgeGroupClassID__in= agegroupclasses,competitionID=kwargs['cmpID'])
        cmpages=[]
        for i in range(0, len(cmpage)):
            cmpages.append(cmpage[i].competitionAgeID)
        schoolclasses = schoolClass.objects.filter(schoolID=kwargs['schoolID'])
        schclasses = []
        for i in range(0, len(schoolclasses)):
            schclasses.append(schoolclasses[i].schoolClassID)
        lists = studentEnrollment.objects.filter(competitionAgeID__in =cmpages,schoolClassID__in=schclasses)
        serializers=studentEnrollmentSerializer(lists,many=True)
        return Response(serializers.data)

class GetStudentsAgeGroupWiseToppers(APIView):                    #Get Toppers AgeGrpWise API
     authentication_classes = (TokenAuthentication, )
     permission_classes = (permissions.IsAuthenticated,)

     def get(self,request,**kwargs):
        agegrpclasses=AgeGroupClass.objects.filter(AgeGroupID=kwargs['AgeID'])
        agegroupclasses=[]
        for i in range(0,len(agegrpclasses)):
            agegroupclasses.append(agegrpclasses[i].AgeGroupClassID)
        cmpage=competitionAge.objects.filter(AgeGroupClassID__in= agegroupclasses,competitionID=kwargs['cmpID'])
        cmpages=[]
        for i in range(0, len(cmpage)):
            cmpages.append(cmpage[i].competitionAgeID)
        schoolclasses = schoolClass.objects.filter(schoolID=kwargs['schoolID'])
        schclasses = []
        for i in range(0, len(schoolclasses)):
            schclasses.append(schoolclasses[i].schoolClassID)
        lists = studentEnrollment.objects.filter(competitionAgeID__in =cmpages,schoolClassID__in=schclasses).exclude(score=999)
        data = []
        lists1 = []
        if len(lists)>0:
            for i in range(0,len(lists)):
                data.append({ 'score': lists[i].score,
                              'time': lists[i].timeTaken,
                              'userID': lists[i].studentEnrollmentID})
            data = sorted(data, key=itemgetter('score', 'time'))
            data = sorted(data, key=lambda k: (-k['score'], k['time']))
            if len(data) >= 3:
                for i in range(0,3):
                    lists1.append(data[i]['userID'])
            else:
                 for i in range(0,len(data)):
                    lists1.append(data[i]['userID'])
            lists = studentEnrollment.objects.filter(studentEnrollmentID__in=lists1)
        serializer = studentEnrollmentSerializer(lists, many=True)
        return Response(serializer.data)

class GetAllStudentsSchoolWise(APIView):                           #Get Students SchoolWise API
     authentication_classes = (TokenAuthentication, )
     permission_classes = (permissions.IsAuthenticated,)

     def get(self,request,**kwargs):
        sch=school.objects.filter(schoolID=kwargs['schoolID'])
        schclasses=schoolClass.objects.filter(schoolID=sch[0].schoolID)
        classes=[]
        for i in range(0,len(schclasses)):
            classes.append(schclasses[i].schoolClassID)
        students=studentEnrollment.objects.filter(schoolClassID__in=classes)
        serializers=studentEnrollmentSerializer(students,many=True)
        return Response(serializers.data)

class UpdateQuestionView(APIView):                         #Update Question API
     authentication_classes = (TokenAuthentication, )
     permission_classes = (permissions.IsAuthenticated,)

     def post(self, request, **kwargs):
        id = correctOption.objects.get(questionTranslationID=kwargs['questionTranslationID'])
        serializers = UpdateQuestions(instance=id, data=request.data, partial=True)
        serializers.is_valid(raise_exception=True)
        serializers.save()
        return Response(serializers.data)

class InsertAgeGrpView(APIView):                         #Insert new AgGrp API
     authentication_classes = (TokenAuthentication, )
     permission_classes = (permissions.IsAuthenticated,)

     def post(self, request, **kwargs):
        serializer = AgeGroupClassSerializer(data=request.data)
        serializer.is_valid(raise_exception=True)
        serializer.save()
        return Response(status=200)

class InsertCompetition(APIView):                          #Insert new Cmp API
     authentication_classes = (TokenAuthentication, )
     permission_classes = (permissions.IsAuthenticated,)

     def post(self, request, **kwargs):
        serializer = CompetitionAgeSerializer(data=request.data)
        serializer.is_valid(raise_exception=True)
        serializer.save()
        return Response(status=200)

class InsertCmpQues(APIView):                             #Insert new CmpQuestion API
     authentication_classes = (TokenAuthentication, )
     permission_classes = (permissions.IsAuthenticated,)

     def post(self, request, **kwargs):
        AgeID = request.data.pop('AgeGroupID',None)
        AgeClassRef = AgeGroupClass.objects.filter(AgeGroupID=AgeID)
        cmpAgeList = []
        for i in range(0,len(AgeClassRef)):
            cmpAgeList.append(AgeClassRef[i].AgeGroupClassID)
        codeRef = code.objects.get(codeName=request.data.pop('cmptype',None))
        CmpRef = competition.objects.get(competitionName = request.data.pop('competitionName',None),startDate = request.data.pop('startdate',None),competitionType =codeRef.codeID )
        CmpAgeRef = competitionAge.objects.filter(AgeGroupClassID__in = cmpAgeList,competitionID=CmpRef)
        quesIDList = request.data.pop('quesList',None)
        for i in range(0,len(quesIDList)):
            LevelcodeRef = code.objects.get(codeName=quesIDList[i]['questionLevelCodeID'])
            for j in range(0,len(CmpAgeRef)):
                t = {
                    "competitionAgeID":CmpAgeRef[j].competitionAgeID,
                    "questionID":quesIDList[i]['questionID'],
                    "questionLevelCodeID":LevelcodeRef.codeID
                }
                serializer = CompetitionQuestionSerializer(data=collections.OrderedDict(t))
                serializer.is_valid(raise_exception=True)
                serializer.save()
        return Response(status=200)

class GetquesAge(APIView):                              #Get Ques from QuestionAge not slected in API
     authentication_classes = (TokenAuthentication, )
     permission_classes = (permissions.IsAuthenticated,)

     def get(self,request,**kwargs):
        ageid = kwargs['AgeID']
        AgeGrp =AgeGroup.objects.get(AgeGroupID =ageid)
        agename = AgeGrp.AgeGroupName.split('-')
        lists = QuestionAge.objects.filter(AgeGroupID=ageid)
        serializers=QuestionAgeSerializer(lists,many=True)
        quesList = []
        for i in range(0,len(lists)):
            quesList.append(lists[i].questionID)
        ageclass = AgeGroupClass.objects.filter(AgeGroupID=ageid)
        ageGrpclass = []
        for i in range(0,len(ageclass)):
            ageGrpclass.append(ageclass[i].AgeGroupClassID)
        cmpage = competitionAge.objects.filter(AgeGroupClassID__in = ageGrpclass,competitionID=kwargs['cmpID'])
        cmpAge = []
        for i in range(0,len(cmpage)):
            cmpAge.append(cmpage[i].competitionAgeID)
        cmpques = competitionQuestion.objects.filter(competitionAgeID__in = cmpAge)
        cmpQues = []
        for i in range(0,len(cmpques)):
            cmpQues.append(cmpques[i].questionID)
        languageRef = code.objects.get(codeName = agename[1])
        quesTransList = questionTranslation.objects.filter(questionID__in =quesList,languageCodeID = languageRef.codeID).exclude(questionID__in = cmpQues)
        serializers1 = GetTranslatedQuestion(quesTransList,many=True)
        return Response({"QuesAge" :serializers.data,
                         "QuesTrans" : serializers1.data})

class InsertOnlyCmpAge(APIView):                             #Insert new agegrp in cmp API
     authentication_classes = (TokenAuthentication, )
     permission_classes = (permissions.IsAuthenticated,)

     def post(self, request, **kwargs):
        serializer = CompetitionAgeOnlySerializer(data=request.data)
        serializer.is_valid(raise_exception=True)
        serializer.save()
        return Response(status=200)


class UpdateCmp(APIView):                         #Update Cmp API
     authentication_classes = (TokenAuthentication, )
     permission_classes = (permissions.IsAuthenticated,)

     def post(self, request, **kwargs):
        cmp_data = request.data.pop('CompetitionData')
        cmp_instance = competition.objects.get(competitionID = request.data.pop('CompetitionID'))
        serializer = CompetitionSerializer(instance=cmp_instance, data=cmp_data, partial=True)
        serializer.is_valid(raise_exception=True)
        serializer.save()
        agedata = request.data.pop('agedata')
        ageRef = AgeGroup.objects.get(AgeGroupID = agedata['ageid'],AgeGroupName=agedata['agename'])
        ageclass = AgeGroupClass.objects.filter(AgeGroupID=ageRef.AgeGroupID)
        ageGrpclass = []
        for i in range(0,len(ageclass)):
            ageGrpclass.append(ageclass[i].AgeGroupClassID)
        cmpage = competitionAge.objects.filter(AgeGroupClassID__in = ageGrpclass,competitionID=cmp_instance.competitionID)
        cmpAge = []
        for i in range(0,len(cmpage)):
            cmpAge.append(cmpage[i].competitionAgeID)
        deleted_data = request.data.pop('DeletedData')
        for i in range(0,len(deleted_data)):
            LevelcodeRef = code.objects.get(codeName=deleted_data[i]['questionLevelCodeID'])
            for j in range(0,len(cmpage)):
                res = competitionQuestion.objects.filter(competitionAgeID = cmpage[j].competitionAgeID,questionID = deleted_data[i]['questionID'],questionLevelCodeID = LevelcodeRef.codeID).exists()
                if res == True:
                    result =competitionQuestion.objects.filter(competitionAgeID = cmpage[j].competitionAgeID,questionID = deleted_data[i]['questionID'],questionLevelCodeID = LevelcodeRef.codeID)
                    result[0].delete()
        cmpQuesList = competitionQuestion.objects.filter(competitionAgeID__in = cmpAge)
        cmpques_data = request.data.pop('CmpQuesData')
        if len(cmpques_data)>0:
            for i in range(0,len(cmpques_data)):
                LevelcodeRef = code.objects.get(codeName=cmpques_data[i]['questionLevelCodeID'])
                quesRef = question.objects.get(questionID = cmpques_data[i]['questionID'])
                for j in range(0,len(cmpage)):
                    competitionQuestion.objects.create(competitionAgeID=cmpage[j],questionID=quesRef,questionLevelCodeID=LevelcodeRef)
        return Response(status=200)


class GetcmpQues(APIView):                       #Get Cmp Question API
     authentication_classes = (TokenAuthentication, )
     permission_classes = (permissions.IsAuthenticated,)

     def get(self,request,**kwargs):
        ageRef = AgeGroup.objects.get(AgeGroupID = kwargs['AgeID'])
        ageclass = AgeGroupClass.objects.filter(AgeGroupID=ageRef.AgeGroupID)
        agename =ageRef.AgeGroupName.split('-')
        ageGrpclass = []
        for i in range(0,len(ageclass)):
            ageGrpclass.append(ageclass[i].AgeGroupClassID)
        cmpage = competitionAge.objects.filter(AgeGroupClassID__in = ageGrpclass,competitionID=kwargs['cmpID'])
        cmpQuesList = competitionQuestion.objects.filter(competitionAgeID=cmpage[0].competitionAgeID)
        serializers = CmpQuesSerializer(cmpQuesList,many=True)
        quesList = []
        for i in range(0,len(cmpQuesList)):
            quesList.append(cmpQuesList[i].questionID)
        languageRef = code.objects.get(codeName = agename[1])
        quesTransList = questionTranslation.objects.filter(questionID__in =quesList,languageCodeID = languageRef.codeID)
        serializers1 = GetTranslatedQuestion(quesTransList,many=True)
        return Response({"cmpQuesList":serializers.data,
                         "QuesTrans":serializers1.data})

class GetClassAgeGrpWise(APIView):                    #Get Classes AgeGrpWise API
     authentication_classes = (TokenAuthentication, )
     permission_classes = (permissions.IsAuthenticated,)

     def get(self,request,**kwargs):
        ageRef = AgeGroup.objects.get(AgeGroupID=kwargs['AgeID'])
        AgeClassList = AgeGroupClass.objects.filter(AgeGroupID=ageRef.AgeGroupID)
        serializers = AgeGroupClassSerializer(AgeClassList,many=True)
        return Response(serializers.data)


class GetAgeGrpCmpWise(APIView):                    #Get AgeGrp  CmpWise API
     authentication_classes = (TokenAuthentication, )
     permission_classes = (permissions.IsAuthenticated,)

     def get(self,request,**kwargs):
        cmpageList = competitionAge.objects.filter(competitionID=kwargs['cmpID']).values('AgeGroupClassID').distinct()
        AgeclassList = []
        for i in range(0,len(cmpageList)):
            AgeclassList.append(cmpageList[i]['AgeGroupClassID'])
        ageclassList = AgeGroupClass.objects.filter(AgeGroupClassID__in = AgeclassList).values('AgeGroupID').distinct()
        AgeGrpList = []
        for i in range(0,len(ageclassList)):
            AgeGrpList.append(ageclassList[i]['AgeGroupID'])
        ageGrpList = AgeGroup.objects.filter(AgeGroupID__in=AgeGrpList)
        serializers = GetAgeGroupsid(ageGrpList,many=True)
        bonus = []
        for i in range(0,len(serializers.data)):
            ageclassRef= AgeGroupClass.objects.filter(AgeGroupID=serializers.data[i]['AgeGroupID'])
            cmpageRef = competitionAge.objects.filter(AgeGroupClassID = ageclassRef[0].AgeGroupClassID,competitionID=kwargs['cmpID'])
            bonus.append(cmpageRef[0].defaultBonusMarks)
        return Response({"AgeGrp":serializers.data,
                         "BonusList":bonus})

class GetMarksAgeWise(APIView):                     #Get AgeGrpWise markingScheme API
     authentication_classes = (TokenAuthentication, )
     permission_classes = (permissions.IsAuthenticated,)

     def get(self,request,**kwargs):
        ageclass = AgeGroupClass.objects.filter(AgeGroupID=kwargs['AgeID'])
        cmpage = competitionAge.objects.filter(AgeGroupClassID=ageclass[0].AgeGroupClassID,competitionID=kwargs['cmpID'])
        Marks = competition_MarkScheme.objects.filter(competitionAgeID=cmpage[0].competitionAgeID)
        serializers = MarkingSchemeSerializer(Marks,many=True)
        return Response(serializers.data)

class UpdateMarkingScheme(APIView):                    #Update markig Scheme API
     authentication_classes = (TokenAuthentication, )
     permission_classes = (permissions.IsAuthenticated,)

     def post(self, request, **kwargs):
         ageclass = AgeGroupClass.objects.filter(AgeGroupID = request.data.pop('AgeID'))
         classID =[]
         for i in range(0,len(ageclass)):
             classID.append(ageclass[i].AgeGroupClassID)
         cmpage = competitionAge.objects.filter(AgeGroupClassID__in = classID,competitionID = request.data.pop('competitionID'))
         levelRef = code.objects.get(codeName = request.data.get('queslevelcode'))
         for i in range(0,len(cmpage)):
             instance = competition_MarkScheme.objects.get(competitionAgeID=cmpage[i].competitionAgeID,questionLevelCodeID=levelRef.codeID)
             t = {
                 "competitionAgeID":cmpage[i].competitionAgeID,
                 "questionLevelCodeID":{
                     "codeName":request.data.get('queslevelcode')
                 },
                 "correctMarks":request.data.get('corrMarks'),
                 "incorrectMarks":request.data.get('incorrMarks')
             }
             serializers = MarkingSchemeSerializer(instance=instance, data=collections.OrderedDict(t), partial=True)
             serializers.is_valid(raise_exception=True)
             serializers.save()
         return Response(status=200)

class GetquesAgeAll(APIView):                       #Get Questionsfrom QuesAge AgeGrpwise API
     authentication_classes = (TokenAuthentication, )
     permission_classes = (permissions.IsAuthenticated,)

     def get(self,request,**kwargs):
        ageid = kwargs['AgeID']
        AgeGrp =AgeGroup.objects.get(AgeGroupID =ageid)
        agename = AgeGrp.AgeGroupName.split('-')
        lists = QuestionAge.objects.filter(AgeGroupID=ageid)
        serializers=QuestionAgeSerializer(lists,many=True)
        quesList = []
        for i in range(0,len(lists)):
            quesList.append(lists[i].questionID)
        languageRef = code.objects.get(codeName = agename[1])
        quesTransList = questionTranslation.objects.filter(questionID__in =quesList,languageCodeID = languageRef.codeID)
        serializers1 = GetTranslatedQuestion(quesTransList,many=True)
        return Response({"QuesAge" :serializers.data,
                         "QuesTrans" : serializers1.data})


class QuesUsage(APIView):                            #Get ques usage API
    authentication_classes = (TokenAuthentication, )
    permission_classes = (permissions.IsAuthenticated,)

    def get(self,request,**kwargs):
        cmpquesList = competitionQuestion.objects.filter(questionID = kwargs['quesID'])
        serializers = CmpQuesGetAllSerializer(cmpquesList,many=True)
        cmpQues = []
        AgegrpName = []
        cmpid = []
        for i in range(0,len(serializers.data)):
            if serializers.data[i]['competitionAgeID']['AgeGroupClassID']['AgeGroupID']['AgeGroupName'] not in AgegrpName or serializers.data[i]['competitionAgeID']['competitionID']['competitionID'] not in cmpid:
                cmpQues.append(serializers.data[i]['competitionQuestionID'])
                AgegrpName.append(serializers.data[i]['competitionAgeID']['AgeGroupClassID']['AgeGroupID']['AgeGroupName'])
                cmpid.append(serializers.data[i]['competitionAgeID']['competitionID']['competitionID'])
        cmpquesList = competitionQuestion.objects.filter(competitionQuestionID__in =cmpQues)
        serializers = CmpQuesGetAllSerializer(cmpquesList,many=True)
        return Response(serializers.data)


class GetAllStudentsAgeGroupWise(APIView):               #Get All Students Age Grp Wise API
     authentication_classes = (TokenAuthentication, )
     permission_classes = (permissions.IsAuthenticated,)

     def get(self,request,**kwargs):
        AgeclassList = AgeGroupClass.objects.filter(AgeGroupID = kwargs['AgeID'])
        AgeClassList = []
        for i in range(0,len(AgeclassList)):
            AgeClassList.append(AgeclassList[i].AgeGroupClassID)
        CmpageList = competitionAge.objects.filter(AgeGroupClassID__in = AgeClassList , competitionID = kwargs['cmpID'])
        CmpAgeList = []
        for i in range(0,len(CmpageList)):
            CmpAgeList.append(CmpageList[i].competitionAgeID)
        lists = studentEnrollment.objects.filter(competitionAgeID__in= CmpAgeList).exclude(score=999)
        serializers=studentEnrollmentSerializer(lists,many=True)
        return Response(serializers.data)

class GetClassWiseAgeGroup(APIView):                  #Get AgeGrps ClassWise API
    authentication_classes = (TokenAuthentication, )
    permission_classes = (permissions.IsAuthenticated,)

    def get(self,request,**kwargs):
        agegrpclass=AgeGroupClass.objects.filter(ClassID=kwargs['Class'])
        a=GetAgeGrpCmpWise()
        d=a.get(request,cmpID=kwargs['cmpID'])
        id = []
        for i in range(0,len(d.data['AgeGrp'])):
            for j in range(0,len(agegrpclass)):
                if(d.data['AgeGrp'][i]['AgeGroupID']==agegrpclass[j].AgeGroupID.AgeGroupID):
                    id.append(agegrpclass[j].AgeGroupID.AgeGroupID)
        agelist = AgeGroup.objects.filter(AgeGroupID__in=id)
        serializers=GetAgeGroupsid(agelist,many=True)
        return Response(serializers.data)

class GetTotalMarks(APIView):                   #Get Total marks API
     authentication_classes = (TokenAuthentication, )
     permission_classes = (permissions.IsAuthenticated,)

     def get(self,request,**kwargs):
        cmp=kwargs['cmpID']
        g=GetcmpQues()
        ques=g.get(request,AgeID=kwargs['ageID'],cmpID=cmp)
        marks=0
        defaultbonus = 0
        for i in range(0,len(ques.data['cmpQuesList'])):
            codeRef = code.objects.get(codeName=ques.data['cmpQuesList'][i]['questionLevelCodeID']['codeName'])
            marks = marks + competition_MarkScheme.objects.filter(competitionAgeID=ques.data['cmpQuesList'][i]['competitionAgeID'],questionLevelCodeID=codeRef.codeID)[0].correctMarks
            defaultbonus = competitionAge.objects.get(competitionAgeID = ques.data['cmpQuesList'][i]['competitionAgeID']).defaultBonusMarks
        total = marks + defaultbonus
        return Response(total)


@parser_classes((MultiPartParser,))
class CustomizePPT(APIView):                      #Get functions related to ppt to zip
    authentication_classes = (TokenAuthentication, )
    permission_classes = (permissions.IsAuthenticated,)

    def saveppt(self,f):
        ppt= Presentation(f)
        ppt.save(os.path.join(settings.MEDIA_ROOT) +'\ppt\\'+ f.name)

    parser_classes = (FileUploadParser,)
    def post(self,request):
        files=request.FILES.getlist('ppt')
        for f in files:
             self.saveppt(f)
        return Response(status=200)

    def ppt(self,f,data,school,type,duplicate):
        print(f)
        text_runs = []
        def duplicate_slide(pres,prs,index):
            template = prs.slides[index]
            try:
                blank_slide_layout = prs.slide_layouts[6]
            except:
                blank_slide_layout = prs.slide_layouts[len(pres.slide_layouts) - 1]

            copied_slide = pres.slides.add_slide(blank_slide_layout)

            for shp in template.shapes:
                el = shp.element
                newel = copy.deepcopy(el)
                copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

            for _, value in six.iteritems(template.part.rels):
                if "notesSlide" not in value.reltype:
                    copied_slide.part.rels.add_relationship(value.reltype,
                                                            value._target,
                                                            value.rId)

            return copied_slide

        i = -1
        while i < (len(data)-1):
            i = i+1
            if type=='participation':
                prs=Presentation(f)
            elif type=='schoolToppers' or type=='nationalToppers':
                prs=Presentation(f)
                if(i == 3):
                    break
                if i==0:
                    for i in range(2, 0, -1):
                        rId = prs.slides._sldIdLst[i].rId
                        prs.part.drop_rel(rId)
                        del prs.slides._sldIdLst[i]
                    i=0
                elif i==1:
                    rId = prs.slides._sldIdLst[2].rId
                    prs.part.drop_rel(rId)
                    del prs.slides._sldIdLst[2]
                    rId = prs.slides._sldIdLst[0].rId
                    prs.part.drop_rel(rId)
                    del prs.slides._sldIdLst[0]
                    i=1
                elif i==2:
                    for i in range(1, -1, -1):
                        rId = prs.slides._sldIdLst[i].rId
                        prs.part.drop_rel(rId)
                        del prs.slides._sldIdLst[i]
                    i=2
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if (shape.text.find('Name')) != -1:
                                    cur_text = run.text
                                    new_text = cur_text.replace(str('Name'), str(data[i]['Name']))
                                    run.text = new_text
                                if (shape.text.find('year')) != -1:
                                    cur_text = run.text
                                    new_text = cur_text.replace(str('year'), str(data[i]['year']))
                                    run.text = new_text
                                if (shape.text.find('class')) != -1:
                                    cur_text = run.text
                                    new_text = cur_text.replace(str('class'), str(data[i]['class']))
                                    run.text = new_text
                                if (shape.text.find('group')) != -1:
                                    cur_text = run.text
                                    new_text = cur_text.replace(str('group'), str(data[i]['group']))
                                    run.text = new_text
                                if (shape.text.find('score')) != -1:
                                    cur_text = run.text
                                    new_text = cur_text.replace(str('score'), str(data[i]['score']))
                                    run.text = new_text
                                if (shape.text.find('total')) != -1:
                                    cur_text = run.text
                                    new_text = cur_text.replace(str('total'), str(data[i]['total']))
                                    run.text = new_text
                                if (shape.text.find('[')) != -1:
                                    cur_text = run.text
                                    new_text = cur_text.replace(str('['), str(''))
                                    run.text = new_text
                                if (shape.text.find(']')) != -1:
                                    cur_text = run.text
                                    new_text = cur_text.replace(str(']'), str(''))
                                    run.text = new_text

            path = ''
            p=''
            if(type=='participation'):
                path=os.path.join(settings.MEDIA_ROOT) + '/output//'+data[i]['Name']+'-'+data[i]['loginID']+'-'+school+'-Class-'+data[0]['class']+'-'+data[0]['group']+'-'+data[0]['year']+'.pptx'
                p=school+'-Class-'+data[0]['class']+'-'+data[0]['group']+'-'+data[0]['year']
                fantasy_zip = zipfile.ZipFile(os.path.join(settings.MEDIA_ROOT) + ("/output/") +school+'-Class-'+data[0]['class']+'-'+data[0]['year']+ '.zip', 'w')
            elif(type=='schoolToppers'):
                path=os.path.join(settings.MEDIA_ROOT) + '\output\\' +data[i]['Name']+'-'+data[i]['loginID']+'-'+ school + '-Toppers-' + data[0][ 'group'] + '-' + data[0]['year'] + '.pptx'
                fantasy_zip = zipfile.ZipFile(os.path.join(settings.MEDIA_ROOT) + ("/output/") + school +'-Toppers-' + data[0][ 'group'] + '-' + data[0]['year'] + '.zip', 'w')
                p=school + '-Toppers-' + data[0][ 'group'] + '-' + data[0]['year']
            elif (type == 'nationalToppers'):
                path = os.path.join(settings.MEDIA_ROOT) + '\output\\' + data[i]['Name']+'-'+data[i]['loginID']+'-'+ 'National Toppers-' + data[0]['group'] + '-' + data[0]['year'] + '.pptx'
                fantasy_zip = zipfile.ZipFile(os.path.join(settings.MEDIA_ROOT) + ("/output/") + 'National Toppers-' + data[0]['group'] + '-' + data[0]['year'] + '.zip', 'w')
                p='National Toppers-' + data[0]['group'] + '-' + data[0]['year']
            prs.save(path)
#            pythoncom.CoInitialize()

            files=glob.glob(path)

            for filename in files:
                command = "unoconv -f pdf '" + filename+"'"
                os.system(command)
                os.remove(path)

#             def convert(files, formatType=32):
#                 for filename in files:
#                     try :
#                         powerpoint = win32com.client.gencache.EnsureDispatch("Powerpoint.Application")
#                         powerpoint.Visible = 1
#                         newname = os.path.splitext(filename)[0] + ".pdf"
#                         deck = powerpoint.Presentations.Open(filename, WithWindow=False)
#                         deck.SaveAs(newname, formatType)
#                         deck.Close()
#                         os.remove(path)
#                     except:
#                         os.remove(path)
#                         os.system('taskkill /F /IM POWERPNT.EXE')
#                 os.system('taskkill /F /IM POWERPNT.EXE')
#             files = glob.glob(path)
#             convert(files)
        f=[]

        for folder, subfolders, files in walk(os.path.join(settings.MEDIA_ROOT) + ("/output/")):
            for file in files:
                if file.endswith('.pdf') and file.find(p)!=-1:
                    fantasy_zip.write(os.path.join(folder, file),
                                      os.path.relpath(os.path.join(folder, file), os.path.join(settings.MEDIA_ROOT) + ("/output/")),
                                      compress_type=zipfile.ZIP_DEFLATED)
                    f.append(file)
        fantasy_zip.close()
        for file in f:
            os.remove(os.path.join(settings.MEDIA_ROOT) + ("/output/")+file)


class deleteFiles(APIView):                       #Delete zip created API
    authentication_classes = (TokenAuthentication, )
    permission_classes = (permissions.IsAuthenticated,)

    c = CustomizePPT()
    ppt = c.ppt
    def get(self,request,**kwargs):
        id=kwargs['schoolID']
        group=kwargs['group']
        year=kwargs['year']
        Class=kwargs['class']
        type=kwargs['CertificateType']
        if(type=='participation'):
            School = school.objects.filter(schoolID=id)
            sch = SchoolSerializers(School, many=True)
            os.remove(os.path.join(settings.MEDIA_ROOT) + '/output//'+sch.data[0]['schoolName']+', '+sch.data[0]['addressID']['city']+'-Class-'+str(Class)+'-'+ year + '.zip')
        if(type=='schoolToppers'):
            School = school.objects.filter(schoolID=id)
            sch = SchoolSerializers(School, many=True)
            os.remove(os.path.join(settings.MEDIA_ROOT) + '\output\\' + sch.data[0]['schoolName'] + ', ' +sch.data[0]['addressID']['city'] + '-Toppers-' +  group + '-' + year + '.zip')
        if(type=='nationalToppers'):
            os.remove(os.path.join(settings.MEDIA_ROOT) + '\output\\' + 'National Toppers-' + group + '-' + year + '.zip')
        return Response(status=200)

class GetParticipationCertificates(APIView):                  #Get Participation Certificate API
    authentication_classes = (TokenAuthentication, )
    permission_classes = (permissions.IsAuthenticated,)

    def get(self, request,**kwargs ):
        c = CustomizePPT()
        g=GetSchoolClassStudents()
        t=GetLatestTemplate()
        print(t)
        d=g.get(request,cmpID=kwargs['cmpID'],schoolClassID=kwargs['schoolClassID'])
        if len(d.data)>0:
            cl=GetClassWiseAgeGroup()
            agegrpid=cl.get(request,cmpID=kwargs['cmpID'], Class=d.data[0]['schoolClassID']['classNumber'])
            data=[]
            template=t.get(request,type='Participation')
            print(template.data)
            for i in range(0,len(d.data)):
                for j in range(0,len(agegrpid.data)):
                    if d.data[i]['competitionAgeID']['AgeGroupClassID']['AgeGroupID']['AgeGroupID']== agegrpid.data[j]['AgeGroupID']:
                         tot=GetTotalMarks()
                         totalMarks=tot.get(request,cmpID=kwargs['cmpID'],ageID=agegrpid.data[j]['AgeGroupID'])
                data.append({'year':d.data[i]['competitionAgeID']['competitionID']['startDate'][0:4],
                             'Name':d.data[i]['userID']['username'],
                             'loginID':d.data[i]['userID']['loginID'],
                             'group':d.data[i]['competitionAgeID']['AgeGroupClassID']['AgeGroupID']['AgeGroupName'],
                             'score':str(d.data[i]['score']+d.data[i]['bonusMarks']),
                             'total':totalMarks.data,
                             'class':str(d.data[i]['schoolClassID']['classNumber'])})
            school=d.data[0]['schoolClassID']['schoolID']['schoolName']+', '+d.data[0]['schoolClassID']['schoolID']['addressID']['city']
            type='participation'
            print(os.path.join(settings.MEDIA_ROOT))
            print(os.path.join(settings.MEDIA_ROOT) + '/ppt//'+template.data)
            c.ppt(os.path.join(settings.MEDIA_ROOT) + '/ppt//'+template.data, data, school,type,duplicate=False)
            return Response(status=200)
        else:
            return Response(status=204)

class GetSchoolToppers(APIView):                     #Get School Toppers Cerificate API
        authentication_classes = (TokenAuthentication, )
        permission_classes = (permissions.IsAuthenticated,)

        def get(self,request,**kwargs):
            c = CustomizePPT()
            t= GetLatestTemplate()
            g=GetStudentsAgeGroupWise()
            d=g.get(request,cmpID=kwargs['cmpID'],AgeID=kwargs['AgeID'],schoolID=kwargs['schoolID'])
            template=t.get(request,type='School_Toppers')
            tot = GetTotalMarks()
            totalMarks = tot.get(request, cmpID=kwargs['cmpID'], ageID=kwargs['AgeID'])
            if len(d.data)>0:
                data=[]
                for i in range(0, len(d.data)):
                    if d.data[i]['score']!= 999:
                        data.append({'year': d.data[i]['competitionAgeID']['competitionID']['startDate'][0:4],
                                     'Name': d.data[i]['userID']['username'],
                                     'loginID':d.data[i]['userID']['loginID'],
                                     'group': d.data[i]['competitionAgeID']['AgeGroupClassID']['AgeGroupID']['AgeGroupName'],
                                     'score': d.data[i]['score']+d.data[i]['bonusMarks'],
                                     'total': totalMarks.data,
                                     'time': d.data[i]['timeTaken'],
                                     'class': str(d.data[i]['schoolClassID']['classNumber'])})
                if len(data)>0:
                    data = sorted(data, key=itemgetter('score', 'time'))
                    data = sorted(data, key=lambda k: (-k['score'], k['time']))
                    school = d.data[0]['schoolClassID']['schoolID']['schoolName'] + ', ' + d.data[0]['schoolClassID']['schoolID']['addressID']['city']
                    type='schoolToppers'
                    c.ppt(os.path.join(settings.MEDIA_ROOT) + '\ppt\\'+template.data, data, school,type,duplicate=False)
                    return Response(status=200)
                else:
                    return Response(status=204)
            else:
                return Response(status=204)


class GetCountryWiseStudents(APIView):                 #Get Students CountryWise API
    authentication_classes = (TokenAuthentication, )
    permission_classes = (permissions.IsAuthenticated,)

    def get(self,request,**kwargs):
        students=studentEnrollment.objects.all().exclude(score=999)
        stu=[]
        for i in range(0,len(students)):
            if(students[i].schoolClassID.schoolID.addressID.countryID.countryID==int(kwargs['countryID'])
            and students[i].competitionAgeID.competitionID.competitionID==int(kwargs['cmpID'])
            and students[i].competitionAgeID.AgeGroupClassID.AgeGroupID.AgeGroupID==int(kwargs['AgeID'])):

                stu.append(students[i].studentEnrollmentID)
        lists = studentEnrollment.objects.filter(studentEnrollmentID__in=stu)
        serializers = studentEnrollmentSerializer(lists, many=True)
        return Response(serializers.data)

class GetCountryWiseToppers(APIView):                            #Get Toppers CountryWise API
    authentication_classes = (TokenAuthentication, )
    permission_classes = (permissions.IsAuthenticated,)

    def get(self,request,**kwargs):
        students=studentEnrollment.objects.all().exclude(score=999)
        stu=[]
        for i in range(0,len(students)):
            if(students[i].schoolClassID.schoolID.addressID.countryID.countryID==int(kwargs['countryID'])
            and students[i].competitionAgeID.competitionID.competitionID==int(kwargs['cmpID'])
            and students[i].competitionAgeID.AgeGroupClassID.AgeGroupID.AgeGroupID==int(kwargs['AgeID'])):
                stu.append(students[i].studentEnrollmentID)
        lists = studentEnrollment.objects.filter(studentEnrollmentID__in=stu)
        data = []
        lists1 = []
        if len(lists)>0:
            for i in range(0,len(lists)):
                data.append({ 'score': lists[i].score,
                              'time': lists[i].timeTaken,
                              'userID': lists[i].studentEnrollmentID})

            data = sorted(data, key=itemgetter('score', 'time'))
            data = sorted(data, key=lambda k: (-k['score'], k['time']))
            if len(data) >= 3:
                for i in range(0,3):
                    lists1.append(data[i]['userID'])
            else:
                 for i in range(0,len(data)):
                    lists1.append(data[i]['userID'])
            lists = studentEnrollment.objects.filter(studentEnrollmentID__in=lists1)
        serializer = studentEnrollmentSerializer(lists, many=True)
        return Response(serializer.data)

class GetStateWiseStudents(APIView):                           #Get Toppers StateWise API
    authentication_classes = (TokenAuthentication, )
    permission_classes = (permissions.IsAuthenticated,)

    def get(self, request, **kwargs):
        students = studentEnrollment.objects.all().exclude(score=999)
        stu = []
        for i in range(0, len(students)):
            if (students[i].schoolClassID.schoolID.addressID.countryID.countryID == int(kwargs['countryID']) and students[i].schoolClassID.schoolID.addressID.stateID.stateID == int(kwargs['stateID'])
                    and students[i].competitionAgeID.competitionID.competitionID == int(kwargs['cmpID'])
                    and students[i].competitionAgeID.AgeGroupClassID.AgeGroupID.AgeGroupID == int(kwargs['AgeID'])):
                stu.append(students[i].studentEnrollmentID)
        lists = studentEnrollment.objects.filter(studentEnrollmentID__in=stu)
        data = []
        lists1 = []
        if len(lists)>0:
            for i in range(0,len(lists)):
                data.append({ 'score': lists[i].score,
                              'time': lists[i].timeTaken,
                              'userID': lists[i].studentEnrollmentID})

            data = sorted(data, key=itemgetter('score', 'time'))
            data = sorted(data, key=lambda k: (-k['score'], k['time']))
            if len(data) >= 3:
                for i in range(0,3):
                    lists1.append(data[i]['userID'])
            else:
                 for i in range(0,len(data)):
                    lists1.append(data[i]['userID'])
            lists = studentEnrollment.objects.filter(studentEnrollmentID__in=lists1)
        serializer = studentEnrollmentSerializer(lists, many=True)
        return Response(serializer.data)

class GetDistrictWiseStudents(APIView):             #Get Toppers DistrictWise API
    authentication_classes = (TokenAuthentication, )
    permission_classes = (permissions.IsAuthenticated,)

    def get(self, request, **kwargs):
        students = studentEnrollment.objects.all().exclude(score=999)
        stu = []
        for i in range(0, len(students)):
            if (students[i].schoolClassID.schoolID.addressID.countryID.countryID == int(kwargs['countryID']) and students[i].schoolClassID.schoolID.addressID.stateID.stateID == int(kwargs['stateID'])
            and students[i].schoolClassID.schoolID.addressID.districtID.districtID == int(kwargs['districtID'])
                    and students[i].competitionAgeID.competitionID.competitionID == int(kwargs['cmpID'])
                    and students[i].competitionAgeID.AgeGroupClassID.AgeGroupID.AgeGroupID == int(kwargs['AgeID'])):
                stu.append(students[i].studentEnrollmentID)
        lists = studentEnrollment.objects.filter(studentEnrollmentID__in=stu)
        data = []
        lists1 = []
        if len(lists)>0:
            for i in range(0,len(lists)):
                data.append({ 'score': lists[i].score,
                              'time': lists[i].timeTaken,
                              'userID': lists[i].studentEnrollmentID})

            data = sorted(data, key=itemgetter('score', 'time'))
            data = sorted(data, key=lambda k: (-k['score'], k['time']))
            if len(data) >= 3:
                for i in range(0,3):
                    lists1.append(data[i]['userID'])
            else:
                 for i in range(0,len(data)):
                    lists1.append(data[i]['userID'])
            lists = studentEnrollment.objects.filter(studentEnrollmentID__in=lists1)
        serializer = studentEnrollmentSerializer(lists, many=True)
        return Response(serializer.data)

class GetSchoolGroupWiseStudents(APIView):                 #Get Toppers SchoolGrpWise API
    authentication_classes = (TokenAuthentication, )
    permission_classes = (permissions.IsAuthenticated,)

    def get(self, request, **kwargs):
        students = studentEnrollment.objects.all().exclude(score=999)
        stu = []
        for i in range(0, len(students)):
            if (students[i].schoolClassID.schoolID.schoolGroupID.codeID == int(kwargs['codeID'])
                and students[i].competitionAgeID.competitionID.competitionID == int(kwargs['cmpID'])
                and students[i].competitionAgeID.AgeGroupClassID.AgeGroupID.AgeGroupID == int(kwargs['AgeID'])):
                stu.append(students[i].studentEnrollmentID)
        lists = studentEnrollment.objects.filter(studentEnrollmentID__in=stu)
        data = []
        lists1 = []
        if len(lists)>0:
            for i in range(0,len(lists)):
                data.append({ 'score': lists[i].score,
                              'time': lists[i].timeTaken,
                              'userID': lists[i].studentEnrollmentID})

            data = sorted(data, key=itemgetter('score', 'time'))
            data = sorted(data, key=lambda k: (-k['score'], k['time']))
            if len(data) >= 3:
                for i in range(0,3):
                    lists1.append(data[i]['userID'])
            else:
                 for i in range(0,len(data)):
                    lists1.append(data[i]['userID'])
            lists = studentEnrollment.objects.filter(studentEnrollmentID__in=lists1)
        serializer = studentEnrollmentSerializer(lists, many=True)
        return Response(serializer.data)

class GetSchoolTypeWiseStudents(APIView):                   #Get Toppers SchoolTypeWise API
    authentication_classes = (TokenAuthentication, )
    permission_classes = (permissions.IsAuthenticated,)

    def get(self, request, **kwargs):
        students = studentEnrollment.objects.all().exclude(score=999)
        stu = []
        for i in range(0, len(students)):
            if (students[i].schoolClassID.schoolID.schoolTypeCodeID.codeID == int(kwargs['codeID'])
                and students[i].competitionAgeID.competitionID.competitionID == int(kwargs['cmpID'])
                and students[i].competitionAgeID.AgeGroupClassID.AgeGroupID.AgeGroupID == int(kwargs['AgeID'])):
                stu.append(students[i].studentEnrollmentID)
        lists = studentEnrollment.objects.filter(studentEnrollmentID__in=stu)
        data = []
        lists1 = []
        if len(lists)>0:
            for i in range(0,len(lists)):
                data.append({ 'score': lists[i].score,
                              'time': lists[i].timeTaken,
                              'userID': lists[i].studentEnrollmentID})

            data = sorted(data, key=itemgetter('score', 'time'))
            data = sorted(data, key=lambda k: (-k['score'], k['time']))
            if len(data) >= 3:
                for i in range(0,3):
                    lists1.append(data[i]['userID'])
            else:
                 for i in range(0,len(data)):
                    lists1.append(data[i]['userID'])
            lists = studentEnrollment.objects.filter(studentEnrollmentID__in=lists1)
        serializer = studentEnrollmentSerializer(lists, many=True)
        return Response(serializer.data)

class GetNationalToppers(APIView):                     #Get National Toppers Certificate API
    authentication_classes = (TokenAuthentication, )
    permission_classes = (permissions.IsAuthenticated,)

    def get(self,request,**kwargs):
        c = CustomizePPT()
        t= GetLatestTemplate()
        g=GetCountryWiseStudents()
        d=g.get(request,cmpID=kwargs['cmpID'],AgeID=kwargs['AgeID'],countryID=kwargs['countryID'])
        tot = GetTotalMarks()
        totalMarks = tot.get(request, cmpID=kwargs['cmpID'], ageID=kwargs['AgeID'])
        template=t.get(request,type='National_Toppers')
        if len(d.data)>0:
            data=[]
            for i in range(0, len(d.data)):
                data.append({'year': d.data[i]['competitionAgeID']['competitionID']['startDate'][0:4],
                             'Name': d.data[i]['userID']['username'],
                             'loginID':d.data[i]['userID']['loginID'],
                             'group': d.data[i]['competitionAgeID']['AgeGroupClassID']['AgeGroupID']['AgeGroupName'],
                             'score': d.data[i]['score']+d.data[i]['bonusMarks'],
                             'total': totalMarks.data,
                             'time': d.data[i]['timeTaken'],
                             'class': str(d.data[i]['schoolClassID']['classNumber'])})

            data = sorted(data, key=itemgetter('score', 'time'))
            data = sorted(data, key=lambda k: (-k['score'], k['time']))
            type='nationalToppers'
            c.ppt(os.path.join(settings.MEDIA_ROOT) + '\ppt\\'+template.data, data, None,type,duplicate=False)
            return Response(status=200)
        else:
         return Response(status=204)

class GetLatestTemplate(APIView):               #Get Latest Template API
    authentication_classes = (TokenAuthentication, )
    permission_classes = (permissions.IsAuthenticated,)

    def get(self,request,**kwargs):
        type=kwargs['type']
        f = []
        for (dirpath, dirnames, filenames) in walk(os.path.join(settings.MEDIA_ROOT) + ("/ppt/")):
            for file in filenames:
                if (file.find(type) != -1):
                    f.append(file)
            break
        datetimes=[]
        for i in range(0,len(f)):
            arr=f[i].split('_')
            datetimes.append(arr[len(arr)-2]+'_'+arr[len(arr)-1].split('.')[0])
            sortedArray = sorted(
            datetimes,
            key=lambda x: datetime.strptime(x, '%Y-%m-%d_%H-%M-%S'), reverse=True
        )
        for i in range(0,len(f)):
            if str(sortedArray[0]) in f[i]:
                name=f[i]
                break
        return Response(name)

class GetStateWiseMean(APIView):             #Get StateWise Mean API
    authentication_classes = (TokenAuthentication, )
    permission_classes = (permissions.IsAuthenticated,)

    def get(self,request,**kwargs):
        ageclassList = AgeGroupClass.objects.filter(AgeGroupID=kwargs['AgeID'])
        AgeclassList = []
        for i in range(0,len(ageclassList)):
            AgeclassList.append(ageclassList[i].AgeGroupClassID)
        cmpageList = competitionAge.objects.filter(AgeGroupClassID__in = AgeclassList,competitionID = kwargs['cmpID'])
        CmpageList = []
        for i in range(0,len(cmpageList)):
            CmpageList.append(cmpageList[i].competitionAgeID)
        studentlist = studentEnrollment.objects.filter(competitionAgeID__in = CmpageList).exclude(score=999)
        stateids = []
        for i in range(0,len(studentlist)):
            if studentlist[i].schoolClassID.schoolID.addressID.stateID.stateID not in stateids:
                stateids.append(studentlist[i].schoolClassID.schoolID.addressID.stateID.stateID)
        stateList = States.objects.filter(stateID__in = stateids)
        FinalList = []
        for i in range(0,len(stateList)):
            marks = []
            for j in range(0,len(studentlist)):
                if studentlist[j].schoolClassID.schoolID.addressID.stateID.stateID == stateList[i].stateID:
                    total = studentlist[j].score+studentlist[j]. bonusMarks
                    marks.append(total)
            if len(marks)!= 0:
                mean = statistics.mean(marks)
            else:
                mean = 0
            stN = States.objects.get(stateID = stateList[i].stateID)
            t = {"state": stN.name,"MeanMarks": mean}
            FinalList.append(t)
        return Response(FinalList)


class CheckAgeGrp(APIView):                 #Check if AgeGrp can update  API
    authentication_classes = (TokenAuthentication, )
    permission_classes = (permissions.IsAuthenticated,)

    def post(self,requests):
        AgeclassRef = AgeGroupClass.objects.filter(AgeGroupID = requests.data.get('AgeID'))
        classesList = []
        classes = []
        for i in range(0,len(AgeclassRef)):
            classesList.append(AgeclassRef[i].AgeGroupClassID)
            classes.append(AgeclassRef[i].ClassID.classNo)
        cmpAgelist = competitionAge.objects.filter(AgeGroupClassID__in=classesList)
        if len(cmpAgelist) >= 1:
            return Response("AgeGroup can't be updated")
        else:
            return Response("Update")

class UpdateAgeGrp(APIView):                       #Get Update Age Grp API
    authentication_classes = (TokenAuthentication, )
    permission_classes = (permissions.IsAuthenticated,)

    def post(self,requests):
            AgeclassRef = AgeGroupClass.objects.filter(AgeGroupID = requests.data.get('AgeID'))
            classesList = []
            classes = []
            for i in range(0,len(AgeclassRef)):
                classesList.append(AgeclassRef[i].AgeGroupClassID)
                classes.append(AgeclassRef[i].ClassID.classNo)
            ageRef = AgeGroup.objects.get(AgeGroupID=requests.data.get('AgeID'))
            t = {
                "AgeGroupName":requests.data.get('AgeGroupName')
            }
            serializer = AgeGroupSerializer(instance=ageRef, data=collections.OrderedDict(t), partial=True)
            serializer.is_valid(raise_exception=True)
            serializer.save()
            lists = requests.data.get('classes')
            for i in range(0,len(AgeclassRef)):
                if AgeclassRef[i].ClassID.classNo not in lists:
                    AgeclassRef[i].delete()
            for i in range(0,len(lists)):
                if lists[i] not in classes:
                    classref = Class.objects.get(classNo=lists[i])
                    AgeGroupClass.objects.create(AgeGroupID = ageRef,ClassID=classref)
            return Response(status=200)


class GetcmpPreview(APIView):                     #Get CmpPreview API
    authentication_classes = (TokenAuthentication, )
    permission_classes = (permissions.IsAuthenticated,)

    def get(self,request,**kwargs):
        ageRef = AgeGroup.objects.get(AgeGroupID = kwargs['AgeID'])
        ageclass = AgeGroupClass.objects.filter(AgeGroupID=ageRef.AgeGroupID)
        agename =ageRef.AgeGroupName.split('-')
        ageGrpclass = []
        for i in range(0,len(ageclass)):
            ageGrpclass.append(ageclass[i].AgeGroupClassID)
        cmpage = competitionAge.objects.filter(AgeGroupClassID__in = ageGrpclass,competitionID=kwargs['cmpID'])
        cmpQuesList = competitionQuestion.objects.filter(competitionAgeID=cmpage[0].competitionAgeID)
        serializers = CmpQuesSerializer(cmpQuesList,many=True)
        quesList = []
        for i in range(0,len(cmpQuesList)):
            quesList.append(cmpQuesList[i].questionID)
        languageRef = code.objects.get(codeName = agename[1])
        quesTransList = questionTranslation.objects.filter(questionID__in =quesList,languageCodeID = languageRef.codeID)
        opts = option.objects.filter(questionID__in=quesList)
        optionList = []
        for i in range(0,len(opts)):
            optionList.append(opts[i].optionID)
        optTrans = optionTranslation.objects.filter(optionID__in=optionList,languageCodeID = languageRef.codeID)
        serializers2 = GetAllTranslatedOptions(optTrans,many=True)
        imagelist = []
        skills=[]
        questransID = []
        for i in range(0,len(quesTransList)):
            if quesTransList[i].questionID.questionTypeCodeID.codeName=='Mcqs_With_Images':
                opt = option.objects.filter(questionID = quesTransList[i].questionID.questionID)
                for j in range(0,len(opt)):
                    coderef = code.objects.get(codeName='ImageOption')
                    oimg = Image.objects.get(ObjectID=opt[j].optionID,ImageTypeCodeID=coderef.codeID)
                    imagelist.append(oimg.ImageID)
            quesref = question.objects.get(questionID = quesTransList[i].questionID.questionID)
            skill = quesref.cs_skills.split(",")
            l1=[]
            for k in range(0,len(skill)):
              l1.append(code.objects.get(codeID=skill[k]).codeName)
            t = {"questionID": quesTransList[i].questionID.questionID,"skills":l1 }
            skills.append(t)
            questransID.append(quesTransList[i].questionTranslationID)
        corrlist = correctOption.objects.filter(questionTranslationID__in =questransID)
        serializers1 = GetCorrectOption(corrlist,many=True)
        Imglists =  Image.objects.filter(ImageID__in = imagelist)
        serializers3 = GetImages(Imglists, many=True)
        return Response({"cmpQuesList":serializers.data,
                         "QuesTrans":serializers1.data,
                         "OptTrans":serializers2.data,
                         "OptImages":serializers3.data,
                         "skills":skills})


class GetAllSchoolStudents(APIView):            #Get All School Students API
    authentication_classes = (TokenAuthentication, )
    permission_classes = (permissions.IsAuthenticated,)

    def get(self,request,**kwargs):
        sch=school.objects.filter(schoolID=kwargs['schoolID'])
        schclasses=schoolClass.objects.filter(schoolID=sch[0].schoolID)
        classes=[]
        for i in range(0,len(schclasses)):
            classes.append(schclasses[i].schoolClassID)
        students=studentEnrollment.objects.filter(schoolClassID__in=classes).values('userID').distinct()
        s=[]
        for i in range(0,len(students)):
            obj=studentEnrollment.objects.filter(userID=students[i]['userID'])[0]
            s.append(obj.studentEnrollmentID)
        students=studentEnrollment.objects.filter(studentEnrollmentID__in =s)
        paginator = CustomPagination()
        response = paginator.generate_response(students,studentEnrollmentSerializer,request)
        userRoleId = []
        for i in range(0,len(response.data['results'])):
            userRoleId.append(response.data['results'][i]['userID']['userID'])
            lists = UserRole.objects.filter(userID__in=userRoleId)
            serializer = UserRoleSerializer(lists,many=True)
        return Response({
            "StudData":response.data,
            "RoleData":serializer.data
        })

class RemoveParticipants(APIView):            #Remove Participants API
    authentication_classes = (TokenAuthentication, )
    permission_classes = (permissions.IsAuthenticated,)

    def post(self,requests):
        studentList = requests.data['studentList']
        for i in range(0,len(studentList)):
            studobj = studentEnrollment.objects.get(studentEnrollmentID=studentList[i])
            studobj.delete()
        return Response(status=200)

class RollbackQuestion(APIView):                 #Rollback process for QuesInsert API
    authentication_classes = (TokenAuthentication, )
    permission_classes = (permissions.IsAuthenticated,)

    def get(self,requests):
        QuesID = question.objects.all().last()
        if questionTranslation.objects.filter(questionID=QuesID.questionID).exists()== True and correctOption.objects.filter(questionTranslationID=questionTranslation.objects.filter(questionID=QuesID.questionID)[0].questionTranslationID).exists()==True:
            return Response(status=200)
        else:
            if questionTranslation.objects.filter(questionID=QuesID.questionID).exists()== False:
                if Image.objects.filter(ObjectID=QuesID.questionID,ImageTypeCodeID=code.objects.get(codeName='ImageQuestion').codeID).exists()==True:
                    imglist = Image.objects.filter(ObjectID=QuesID.questionID,ImageTypeCodeID=code.objects.get(codeName='ImageQuestion').codeID)
                    for i in range(0,len(imglist)):
                        imglist[i].delete()
                if Image.objects.filter(ObjectID=QuesID.questionID,ImageTypeCodeID=code.objects.get(codeName='ImageAnsExplanation').codeID).exists()==True:
                    imglist = Image.objects.filter(ObjectID=QuesID.questionID,ImageTypeCodeID=code.objects.get(codeName='ImageAnsExplanation').codeID)
                    for i in range(0,len(imglist)):
                        imglist[i].delete()
                if QuestionAge.objects.filter(questionID=QuesID.questionID).exists()==True:
                    agelist = QuestionAge.objects.filter(questionID=QuesID.questionID)
                    for i in range(0,len(agelist)):
                        agelist[i].delete()
                QuesID.delete()
            else:
                 if questionTranslation.objects.filter(questionID=QuesID.questionID).exists()== True and correctOption.objects.filter(questionTranslationID=questionTranslation.objects.filter(questionID=QuesID.questionID)[0].questionTranslationID).exists()==False:
                     if Image.objects.filter(ObjectID=QuesID.questionID, ImageTypeCodeID=code.objects.get(
                             codeName='ImageQuestion').codeID).exists() == True:
                         imglist = Image.objects.filter(ObjectID=QuesID.questionID, ImageTypeCodeID=code.objects.get(
                             codeName='ImageQuestion').codeID)
                         for i in range(0, len(imglist)):
                             imglist[i].delete()
                     if Image.objects.filter(ObjectID=QuesID.questionID, ImageTypeCodeID=code.objects.get(
                             codeName='ImageAnsExplanation').codeID).exists() == True:
                         imglist = Image.objects.filter(ObjectID=QuesID.questionID, ImageTypeCodeID=code.objects.get(
                             codeName='ImageAnsExplanation').codeID)
                         for i in range(0, len(imglist)):
                             imglist[i].delete()
                     if QuestionAge.objects.filter(questionID=QuesID.questionID).exists() == True:
                         agelist = QuestionAge.objects.filter(questionID=QuesID.questionID)
                         for i in range(0, len(agelist)):
                             agelist[i].delete()
                     quesTrans = questionTranslation.objects.filter(questionID=QuesID.questionID)
                     quesTrans[0].delete()
                     if option.objects.filter(questionID = QuesID.questionID).exists()==True:
                         optlist = option.objects.filter(questionID = QuesID.questionID)
                         for i in range(0,len(optlist)):
                            if optionTranslation.objects.filter(optionID=optlist[i].optionID).exists() == True:
                                optionTranslation.objects.filter(optionID=optlist[i].optionID)[0].delete()
                            if Image.objects.filter(ObjectID=optlist[i].optionID,ImageTypeCodeID=code.objects.get(codeName='ImageOption').codeID).exists()==True:
                                Image.objects.filter(ObjectID=optlist[i].optionID,ImageTypeCodeID=code.objects.get(codeName='ImageOption').codeID)[0].delete()
                            optlist[i].delete()
                     QuesID.delete()
        return Response(status=200)

class RollbackQuestionTranslation(APIView):  # Rollback process for Translation API
        authentication_classes = (TokenAuthentication, )
        permission_classes = (permissions.IsAuthenticated,)

        def get(self, requests):
            quesTransID = questionTranslation.objects.all().last()
            if correctOption.objects.filter(questionTranslationID=quesTransID.questionTranslationID).exists() == True:
                return Response(status=200)
            else:
                optlist = option.objects.filter(questionID=quesTransID.questionID)
                language = quesTransID.languageCodeID
                for i in range(0, len(optlist)):
                    if optionTranslation.objects.filter(optionID=optlist[i].optionID,
                                                        languageCodeID=language).exists() == True:
                        optionTranslation.objects.filter(optionID=optlist[i].optionID, languageCodeID=language)[
                            0].delete()
                quesTransID.delete()
            return Response(status=200)
